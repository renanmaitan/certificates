import comtypes.client
from pptx import Presentation

def replace_placeholders(caminho_arquivo, nome, cpf, caminho_saida):
    prs = Presentation(caminho_arquivo)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragrafo in shape.text_frame.paragraphs:
                    for run in paragrafo.runs:
                        if "{NOME}" in run.text or "{CPF}" in run.text:
                            run.text = run.text.replace("{NOME}", nome).replace("{CPF}", cpf)

    prs.save(caminho_saida)
    
def save_as_pdf(caminho_pptx, caminho_pdf):
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1

    presentation = powerpoint.Presentations.Open(caminho_pptx, WithWindow=False)

    presentation.SaveAs(caminho_pdf, FileFormat=32)
    presentation.Close()
    powerpoint.Quit()